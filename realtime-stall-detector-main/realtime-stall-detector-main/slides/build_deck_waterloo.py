"""Rebuild the pitch deck on the official UWaterloo template (slides/waterloo_base.pptx).

waterloo_base.pptx is the school .potx with its content-type flipped to a presentation so
python-pptx can open it (see the conversion step in the chat / build notes). We strip the
template's sample slides and rebuild our 8 slides using its layouts, so the deck carries
Waterloo branding (fonts, colours, crest). Real charts live in slides/assets and plots/.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
A = os.path.join(HERE, "assets")
REPO = os.path.dirname(HERE)

prs = Presentation(os.path.join(HERE, "waterloo_base.pptx"))

# strip the template's sample slides (drop the relationship so the part isn't re-serialized)
ids = prs.slides._sldIdLst
for sid in list(ids):
    prs.part.drop_rel(sid.get(qn("r:id")))
    ids.remove(sid)

L = {lay.name: lay for lay in prs.slide_layouts}
TITLE = L["Title Slide"]
CONTENT = L["Title and Content"]
CLOSING = L["Closing Slide"]


def body_bullets(slide, bullets, idx=1):
    tf = slide.placeholders[idx].text_frame
    tf.clear(); tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b; p.space_after = Pt(10)
    return slide.placeholders[idx]


def left_text_right_img(slide, bullets, img, extra=None):
    bp = body_bullets(slide, bullets)
    bp.left = Inches(0.55); bp.top = Inches(1.9); bp.width = Inches(6.3); bp.height = Inches(4.8)
    if extra:
        for line, color in extra:
            p = bp.text_frame.add_paragraph(); r = p.add_run(); r.text = line
            r.font.size = Pt(15); r.font.color.rgb = color; p.space_after = Pt(2)
    slide.shapes.add_picture(img, Inches(7.15), Inches(2.0), width=Inches(5.7))


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


from pptx.dml.color import RGBColor
RED = RGBColor(0xC8, 0x10, 0x2E); ORANGE = RGBColor(0xC2, 0x82, 0x00)

# 1 · Title
s = prs.slides.add_slide(TITLE)
s.placeholders[0].text = "Reading a motor's current to catch failure"
s.placeholders[1].text = "Real-time stall detection — live, and on the chip\nToyota Innovation Challenge · Fault Prediction"
notes(s, "Unexpected machine downtime costs manufacturers millions a year. The earliest warning a motor "
         "gives before it fails is its own current — so we read that current in real time and catch a stall "
         "the instant it happens.")

# 2 · Problem + insight
s = prs.slides.add_slide(CONTENT)
s.placeholders[0].text = "The trap: a stall looks like a start-up"
left_text_right_img(s, [
    "Unplanned downtime costs manufacturers millions / year.",
    "A motor's earliest warning is its own current.",
    "A stall spikes the current — but so does every start-up (inrush).",
    "The signal isn't high current. It's high current that STAYS high.",
    "Stall = sustained. Start-up = decays. Telling them apart is the whole game.",
], os.path.join(A, "signatures_real.png"))
notes(s, "When a motor stalls, current spikes — but it also spikes on every start-up, the inrush. A naive "
         "'too-high' alarm cries wolf on every power-on. The real signal is high current that STAYS high. A "
         "stall is sustained; a start-up decays. Telling them apart is the heart of the project.")

# 3 · How it works
s = prs.slides.add_slide(CONTENT)
s.placeholders[0].text = "Current → 6 features → an interpretable tree"
bp = body_bullets(s, [
    "Slice the current into short windows; extract 6 features — level, spread, slope, trend.",
    "Feed a small decision tree — interpretable, and tiny enough to run on a chip.",
    "A debounced state machine turns it into green / orange / red.",
    "The rule it LEARNED (not hand-written):",
])
for line, color in [("    mean high  &  steady (low spread)  →  STALL", RED),
                    ("    mean high  &  spiky / decaying      →  start-up  (no alarm)", ORANGE)]:
    p = bp.text_frame.add_paragraph(); r = p.add_run(); r.text = line
    r.font.bold = True; r.font.color.rgb = color
notes(s, "Six windowed features — level, slope, and whether the current is still rising or holding — into a "
         "decision tree. We chose a tree: interpretable, and tiny enough to run on a chip. The algorithm "
         "itself found that steadiness separates a real stall from a start-up.")

# 4 · Live demo
s = prs.slides.add_slide(CONTENT)
s.placeholders[0].text = "Live demo"
left_text_right_img(s, [
    "Running  →  GREEN.",
    "Power-cycle a few times  →  stays green (no false alarm on inrush).",
    "Hold the shaft (stall)  →  RED in ~0.1 s.",
    "Release  →  back to green.",
], os.path.join(REPO, "plots", "dashboard_snapshot.png"))
notes(s, "Running — green. Power-cycle it — current spikes but it stays green, no false alarm. Stall it — RED "
         "in under a tenth of a second. Release — back to green. Start-up ignored, real stall caught instantly.")

# 5 · Results
s = prs.slides.add_slide(CONTENT)
s.placeholders[0].text = "It works — and never cries wolf"
left_text_right_img(s, [
    "Stall: 100% separated — never missed, never falsely raised.",
    "Start-up is NEVER predicted as a stall (the red-alarm property).",
    "18 automated tests; a physics simulator verified it before any hardware.",
    "The tree is readable — we can show a judge exactly why it fired.",
], os.path.join(A, "confusion_real.png"))
notes(s, "It cleanly separates the three states. Stall is perfect, and a start-up is never predicted as a "
         "stall — the property that matters. 18 automated tests, simulator-verified before hardware.")

# 6 · Edge ML
s = prs.slides.add_slide(CONTENT)
s.placeholders[0].text = "No laptop: the tree runs on the chip"
body_bullets(s, [
    "We compiled the same decision tree into C and flashed it onto the Arduino itself.",
    "The board runs inference on-chip and lights an LED on a stall.",
    "Pull the laptop — power it from a charger — and it still works. No PC, no cloud.",
    "Detection at the edge, on a microcontroller — where a factory sensor lives.",
])
notes(s, "It doesn't need a laptop. We compiled the same tree into C and flashed it onto the Arduino; it runs "
         "inference on-chip and lights an LED on a stall. Pull the laptop and it still fires — at the edge.")

# 7 · Scale (option B footnote)
s = prs.slides.add_slide(CONTENT)
s.placeholders[0].text = "The method scales to the factory floor"
left_text_right_img(s, [
    "The exact same feature code runs unchanged on Toyota's real 8-DOF arm telemetry — fleet-wide.",
    "An unsupervised pass flags the most mechanically-stressed joints (shoulder/elbow) — physically right.",
    "With labeled fault data, the same approach becomes production predictive maintenance.",
], os.path.join(REPO, "toyota_transfer", "plots", "fleet_anomaly_heatmap.png"),
    extra=[("Real faults aren't labeled here — this shows the method transfers, not arm fault detection.",
            RGBColor(0x70, 0x70, 0x70))])
notes(s, "The same feature code runs unchanged on Toyota's real 8-DOF arm telemetry, fleet-wide, and surfaces "
         "the most stressed joints. With labeled fault data it scales to production predictive maintenance. "
         "(Honest: this shows the method transfers — not that we detected stalls on the arm; that data is "
         "unlabeled.)")

# 8 · Closing (the template's branded closing layout — keep text in the title placeholder)
s = prs.slides.add_slide(CLOSING)
s.placeholders[0].text = ("From a $3 motor to a robot fleet — same physics, same code.   "
                          "We turn a motor's current into an early-warning system.   Thank you.")
notes(s, "It's human-centered: catches failures before they cascade — protecting equipment, preventing line "
         "stops, keeping workers safe. From a $3 motor to a robot fleet — same physics, same code. Thank you.")

out = os.path.join(HERE, "pitch_waterloo.pptx")
prs.save(out)
print("saved", out, "| slides:", len(list(prs.slides)))
