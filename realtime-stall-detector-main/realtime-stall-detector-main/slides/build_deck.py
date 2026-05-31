"""Build the pitch deck (slides/pitch.pptx) with python-pptx. Speaker script lives in the notes."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
A = os.path.join(HERE, "assets")

BG = RGBColor(0x0E, 0x11, 0x16)
PANEL = RGBColor(0x17, 0x1C, 0x24)
TXT = RGBColor(0xE6, 0xED, 0xF3)
MUT = RGBColor(0x9A, 0xA5, 0xB1)
GREEN = RGBColor(0x2E, 0xA0, 0x43)
ORANGE = RGBColor(0xD2, 0x99, 0x22)
RED = RGBColor(0xDA, 0x36, 0x33)
BLUE = RGBColor(0x58, 0xA6, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    return tf


def para(tf, text, size, color=TXT, bold=False, first=False, bullet=False, space=8, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space)
    run = p.add_run(); run.text = ("•  " if bullet else "") + text
    run.font.size = Pt(size); run.font.color.rgb = color; run.font.bold = bold
    run.font.name = "Helvetica Neue"
    return p


def accent(s, color, t=Inches(1.18)):
    bar = s.shapes.add_shape(1, Inches(0.7), t, Inches(2.4), Pt(4))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def img(s, path, l, t, w):
    return s.shapes.add_picture(path, Inches(l), Inches(t), width=Inches(w))


# ---------- 1 · Title ----------
s = slide()
tf = box(s, 0.9, 2.5, 11.5, 2.6)
para(tf, "Reading a motor's current to catch failure", 40, TXT, True, first=True)
para(tf, "Real-time stall detection — live, and on the chip", 26, BLUE, space=24)
para(tf, "Toyota Innovation Challenge  ·  Fault Prediction", 18, MUT)
bar = s.shapes.add_shape(1, Inches(0.95), Inches(2.45), Inches(0.12), Inches(2.0))
bar.fill.solid(); bar.fill.fore_color.rgb = GREEN; bar.line.fill.background()
notes(s, "Unexpected machine downtime costs manufacturers millions a year. The earliest warning a "
         "motor gives before it fails is its own current — so we built a system that reads that current "
         "in real time and catches a stall the instant it happens.")

# ---------- 2 · Problem + insight ----------
s = slide(); accent(s, ORANGE)
para(box(s, 0.7, 0.45, 12, 1), "The trap: a stall looks like a start-up", 32, TXT, True, first=True)
tf = box(s, 0.7, 1.6, 6.0, 5)
for b in ["Unplanned downtime costs manufacturers millions / year.",
          "A motor's earliest warning is its own current.",
          "A stall spikes the current — but so does every start-up (inrush).",
          "The signal isn't *high* current. It's high current that *stays* high.",
          "Stall = sustained.  Start-up = decays.  Telling them apart is the whole game."]:
    para(tf, b, 19, TXT, bullet=True, space=14)
img(s, os.path.join(A, "signatures_real.png"), 7.0, 2.0, 6.0)
notes(s, "Here's what makes this hard. When a motor stalls, current spikes — but current also spikes "
         "every time the motor starts, the inrush. A naive 'current too high' alarm cries wolf on every "
         "power-on. The real signal isn't high current — it's high current that STAYS high. A stall is "
         "sustained; a start-up just decays. Telling those apart is the heart of our project.")

# ---------- 3 · How it works ----------
s = slide(); accent(s, BLUE)
para(box(s, 0.7, 0.45, 12, 1), "Current → 6 features → an interpretable tree", 32, TXT, True, first=True)
tf = box(s, 0.7, 1.6, 11.9, 5)
for b in ["Slice the current into short windows; extract 6 features — level, spread, slope, trend.",
          "Feed a small decision tree — chosen on purpose: interpretable, and tiny enough to run on a chip.",
          "A debounced state machine turns it into green / orange / red.",
          ]:
    para(tf, b, 20, TXT, bullet=True, space=14)
tf2 = box(s, 0.7, 4.4, 11.9, 2.3)
para(tf2, "The rule it LEARNED (not hand-written):", 18, MUT, True, first=True, space=8)
para(tf2, "mean high  &  steady (low spread)  →  STALL", 22, RED, True, space=4)
para(tf2, "mean high  &  spiky / decaying      →  start-up  (no alarm)", 22, ORANGE, True)
notes(s, "Under the hood: six windowed features — level, slope, and whether the current is still rising or "
         "holding. We feed a decision tree, on purpose: it's interpretable — here are the exact rules it "
         "learned — and it's tiny. The algorithm itself discovered that steadiness (low spread) is what "
         "separates a real stall from a start-up inrush.")

# ---------- 4 · Live demo ----------
s = slide(); accent(s, GREEN)
para(box(s, 0.7, 0.45, 12, 1), "Live demo", 32, TXT, True, first=True)
tf = box(s, 0.7, 1.6, 5.6, 5)
para(tf, "running", 24, GREEN, True, first=True, space=6); para(tf, "→ GREEN", 20, MUT, space=18)
para(tf, "power-cycle a few times", 24, ORANGE, True, space=6); para(tf, "→ stays green — no false alarm on inrush", 20, MUT, space=18)
para(tf, "hold the shaft (stall)", 24, RED, True, space=6); para(tf, "→ RED in ~0.1 s", 20, MUT)
img(s, os.path.join(A, "..", "..", "plots", "dashboard_snapshot.png"), 6.5, 2.2, 6.3)
notes(s, "Let me show you. Running — green. Watch when I power-cycle it — current spikes, but it stays "
         "GREEN, no false alarm. Now I actually stall it — RED, in under a tenth of a second. Release — "
         "back to green. Start-up ignored, real stall caught instantly.")

# ---------- 5 · Results ----------
s = slide(); accent(s, GREEN)
para(box(s, 0.7, 0.45, 12, 1), "It works — and never cries wolf", 32, TXT, True, first=True)
tf = box(s, 0.7, 1.7, 6.2, 5)
for b in ["Stall: 100% separated — never missed, never falsely raised.",
          "Start-up is NEVER predicted as a stall (the red-alarm property).",
          "18 automated tests; a physics simulator verified the pipeline before any hardware.",
          "The tree is readable — we can show a judge exactly why it fired."]:
    para(tf, b, 19, TXT, bullet=True, space=14)
img(s, os.path.join(A, "confusion_real.png"), 7.3, 1.7, 5.4)
notes(s, "On our collected data the detector cleanly separates the three states. Stall is perfectly "
         "classified, and crucially a start-up is never predicted as a stall — that's the property that "
         "matters. 18 automated tests, simulator-verified before we touched hardware.")

# ---------- 6 · Edge ML ----------
s = slide(); accent(s, BLUE)
para(box(s, 0.7, 0.45, 12, 1), "No laptop: the tree runs on the chip", 32, TXT, True, first=True)
tf = box(s, 0.7, 1.8, 11.9, 5)
for b in ["We compiled the same decision tree into C and flashed it onto the Arduino itself.",
          "The board runs inference on-chip and lights an LED on a stall.",
          "Pull the laptop — power it from a charger — and it still works. No PC, no cloud.",
          "Detection at the edge, on a microcontroller — exactly where a factory sensor lives."]:
    para(tf, b, 21, TXT, bullet=True, space=16)
notes(s, "And it doesn't need a laptop. We compiled that same tree into C and flashed it onto the Arduino. "
         "The board runs inference on-chip and lights this LED on a stall. Pull the laptop and it still "
         "fires — detection at the edge, on a microcontroller.")

# ---------- 7 · Scale (Toyota footnote, option B) ----------
s = slide(); accent(s, ORANGE)
para(box(s, 0.7, 0.45, 12, 1), "The method scales to the factory floor", 32, TXT, True, first=True)
tf = box(s, 0.7, 1.7, 6.2, 5)
for b in ["The exact same feature code runs unchanged on Toyota's real 8-DOF arm telemetry — fleet-wide.",
          "An unsupervised pass flags the most mechanically-stressed joints (the shoulder/elbow) — physically right.",
          "With labeled fault data, the same approach becomes production predictive maintenance."]:
    para(tf, b, 18, TXT, bullet=True, space=14)
para(tf, "(Real faults aren't labeled in this data — this shows the method transfers, not arm fault detection.)",
     14, MUT, space=4)
img(s, os.path.join(A, "..", "..", "toyota_transfer", "plots", "fleet_anomaly_heatmap.png"), 7.0, 2.1, 6.1)
notes(s, "And the method isn't locked to our bench — the same feature code runs unchanged on Toyota's real "
         "8-DOF arm telemetry, across the whole fleet, and surfaces the most stressed joints. With labeled "
         "fault data it scales to production predictive maintenance. (Be honest: this shows the method "
         "transfers — we are NOT claiming we detected stalls on the arm; that data is unlabeled.)")

# ---------- 8 · Safety + Human + close ----------
s = slide(); accent(s, GREEN)
para(box(s, 0.7, 0.45, 12, 1), "Catches failures before they cascade", 32, TXT, True, first=True)
tf = box(s, 0.7, 1.8, 11.9, 4)
for b in ["Safety: an early stall alarm prevents motor burnout and equipment damage.",
          "Human-centered: less unplanned downtime, fewer line stops, workers kept out of harm's way.",
          "Built and verified safely on a low-voltage, current-limited bench."]:
    para(tf, b, 20, TXT, bullet=True, space=16)
tf2 = box(s, 0.7, 5.2, 11.9, 1.6)
para(tf2, "From a $3 motor to a robot fleet — same physics, same code.", 24, BLUE, True, first=True, space=6)
para(tf2, "We turn a motor's current into an early-warning system.   Thank you.", 22, TXT, True)
notes(s, "It's human-centered: it catches failures before they cascade — protecting equipment, preventing "
         "unplanned line stops, keeping workers out of harm's way. From a 3-dollar motor to a robot fleet — "
         "same physics, same code. We turn a motor's current into an early-warning system. Thank you.")

out = os.path.join(HERE, "pitch.pptx")
prs.save(out)
print("saved", out, "—", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
